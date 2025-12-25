# routes/pptgen.py
import io
import os
import time
import requests
import markdown
from flask import Blueprint, request, jsonify, send_file, current_app
from pptx import Presentation
from pptx.util import Pt, Inches

# load environment variables earlier in your app (e.g. with python-dotenv)
ppt_bp = Blueprint("pptgen", __name__, url_prefix="/ppt")

# ---------- Config ----------
BATCH_SIZE = 6         # how many slides to generate per LLM call
PROMPT_TIMEOUT = 12    # seconds for model HTTP calls
MAX_SLIDES = 40        # guardrails

# Order of GROQ fallback models (try most capable / highest quota first)
GROQ_MODEL_ORDER = [
    "llama-3.1-8b-instant",
    "llama-3.3-70b-versatile",
    "meta-llama/llama-4-scout-17b-16e-instruct",
    "openai/gpt-oss-20b",
    "qwen/qwen3-32b",
    "gemma2-9b-it",
    "meta-llama/llama-4-maverick-17b-128e-instruct"
]


# ---------- Helper: call remote model (OpenRouter -> GROQ models) ----------
def call_openrouter(openrouter_key, model, messages, timeout=PROMPT_TIMEOUT):
    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {openrouter_key}",
        "Content-Type": "application/json",
        "HTTP-Referer": "http://localhost",
        "X-Title": "Mentora PPT Generator",
    }
    payload = {"model": model, "messages": messages}
    resp = requests.post(url, headers=headers, json=payload, timeout=timeout)
    resp.raise_for_status()
    return resp.json()

def call_groq(groq_key, model, messages, timeout=PROMPT_TIMEOUT):
    url = "https://api.groq.com/openai/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {groq_key}",
        "Content-Type": "application/json",
    }
    payload = {"model": model, "messages": messages}
    resp = requests.post(url, headers=headers, json=payload, timeout=timeout)
    resp.raise_for_status()
    return resp.json()

def call_model_with_fallback(openrouter_key, groq_key, prefer_openrouter_model, prefer_groq_models, messages):
    # 1) try openrouter
    if openrouter_key:
        try:
            data = call_openrouter(openrouter_key, prefer_openrouter_model, messages)
            # detect error field
            if isinstance(data, dict) and "error" in data:
                raise Exception(f"OpenRouter error: {data['error']}")
            return data
        except Exception as e:
            current_app.logger.info(f"OpenRouter failed: {e}")

    # 2) try GROQ models in order
    if groq_key:
        for gm in prefer_groq_models:
            try:
                data = call_groq(groq_key, gm, messages)
                if isinstance(data, dict) and "error" in data:
                    raise Exception(f"GROQ error on {gm}: {data['error']}")
                return data
            except Exception as e:
                current_app.logger.info(f"GROQ model {gm} failed: {e}")
                # try next one
    raise RuntimeError("All model providers failed")


# ---------- Prompt building (improved) ----------
BASE_BATCH_PROMPT = """
You MUST output slides ONLY in this strict format. DO NOT use Markdown or any other style.

Example of EXACT required format:

# Slide 1
Title: Neural Networks Overview
Paragraph: This slide explains the basics of neural networks in 3–4 sentences.
Bullets:
- Bullet one with 10–14 words.
- Bullet two with 10–14 words.
- Bullet three with 10–14 words.

-------------------------------------

Now generate {count} slides for TOPIC: "{topic}"

Start numbering from Slide {start_index}.
Use this EXACT format:

# Slide {start_index}
Title: <short title>
Paragraph: <3–4 sentence paragraph>
Bullets:
- <bullet 1>
- <bullet 2>
- <bullet 3>

DESCRIPTION (USE FOR TONE, NOT CONTENT):
{description}

STRICT RULES:
- DO NOT write "# Slide X: Title". Title MUST be on "Title:" line ONLY.
- DO NOT skip "Paragraph:" label.
- DO NOT skip "Bullets:" label.
- DO NOT place bullets before the paragraph.
- DO NOT output extra sentences or Markdown.
- NO bold (** **)
- NO backticks
- NO JSON
- ONLY the exact required format.
"""
# ---------- Parse model text into slides ----------
def parse_slides(text):
    """
    Parse text returned by model into slide dicts:
    [{'title': '...', 'bullets': ['a', 'b', 'c']}, ...]
    """
    slides = []
    if not text:
        return slides

    # Normalize CRLF
    text = text.replace("\r\n", "\n")

    # Split by '# Slide' occurrences (robust)
    parts = []
    for chunk in text.split("\n# Slide"):
        chunk = chunk.strip()
        if chunk:
            # ensure starts with 'Slide' or 'Slide X'
            if chunk.lower().startswith("slide") or chunk.startswith("Slide"):
                parts.append(chunk)
            else:
                # If initial chunk doesn't start with 'Slide', the model may not include '# ' for first slide
                if chunk.startswith("# Slide") or chunk.startswith("Slide"):
                    parts.append(chunk)
                else:
                    # also accept blocks that begin with "Slide X" without '#'
                    if chunk.lower().startswith("slide"):
                        parts.append(chunk)

    for block in parts:
        lines = [l.strip() for l in block.split("\n") if l.strip()]
        # Find title line

        title = None
        paragraph = ""
        bullets = []

        for i, ln in enumerate(lines):
            if ln.lower().startswith("title:"):
                title = ln.split(":", 1)[1].strip()

            elif ln.lower().startswith("paragraph:"):
                paragraph = ln.split(":", 1)[1].strip()

            elif ln.startswith("-"):
                bullets.append(ln[1:].strip())



        if title and len(bullets) >= 3:
            # take exactly first 3 bullets
            slides.append({
                "title": title,
                "paragraph": paragraph,
                "bullets": bullets[:3]
            })


    return slides

# ---------- PPT generation in-memory ----------
def build_ppt_bytes(topic, slides, config=None):
    """
    Return PPTX file bytes (in-memory) for download.
    """
    prs = Presentation()

    # sizes
    TITLE_FONT_SIZE = Pt(40)
    SLIDE_TITLE_SIZE = Pt(28)
    BULLET_FONT_SIZE = Pt(18)
    BULLET_LINE_SPACING = 1.1

    # -------- TITLE SLIDE --------
    slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(slide_layout)

    # Main Title
    title_slide.shapes.title.text = topic
    try:
        tf = title_slide.shapes.title.text_frame.paragraphs[0]
        tf.font.size = TITLE_FONT_SIZE
        tf.font.bold = True
    except Exception:
        pass

    # Subtitle (NO LONGER USING DESCRIPTION)
    try:
        if title_slide.placeholders and len(title_slide.placeholders) > 1:
            subtitle = title_slide.placeholders[1]
            subtitle.text = "Powered by Mentora AI"   # <-- optional clean subtitle
            # subtitle.text = ""                     # <-- OR remove subtitle fully
    except Exception:
        pass

    # -------- CONTENT SLIDES --------
    for idx, s in enumerate(slides, start=1):
        layout = prs.slide_layouts[1]
        sld = prs.slides.add_slide(layout)

        # Title
        try:
            title_tf = sld.shapes.title.text_frame
            title_tf.text = s["title"]
            title_tf.paragraphs[0].font.size = SLIDE_TITLE_SIZE
            title_tf.paragraphs[0].font.bold = True
        except Exception:
            pass

        # Body
        try:
            body = sld.placeholders[1].text_frame
            body.clear()

            # Paragraph (NEW)
            if s.get("paragraph"):
                p = body.add_paragraph()
                p.text = s["paragraph"]
                p.font.size = BULLET_FONT_SIZE
                p.level = 0
                p.line_spacing = BULLET_LINE_SPACING

            # Bullets
            for b in s["bullets"]:
                q = body.add_paragraph()
                q.text = b
                q.level = 1
                q.font.size = BULLET_FONT_SIZE
                q.line_spacing = BULLET_LINE_SPACING
        except Exception:
            pass

    # Save to bytes
    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio

# ---------- Main route ----------
@ppt_bp.route("/generate", methods=["POST"])
def generate_ppt_route():
    """
    POST JSON expected:
    {
      "topic": "The Future of Renewable Energy",
      "description": "target audience and tone (optional)",
      "slides": 12
    }
    Returns: file download (pptx) or JSON error
    """
    payload = request.get_json(force=True, silent=True) or request.form.to_dict()
    topic = payload.get("topic") or payload.get("title")
    description = payload.get("description", "") or ""
    try:
        slides_requested = int(payload.get("slides", 10))
    except Exception:
        slides_requested = 10

    # basic validation
    if not topic or slides_requested < 1 or slides_requested > MAX_SLIDES:
        return jsonify({"error": "Invalid input. Provide 'topic' and slides between 1 and %d." % MAX_SLIDES}), 400

    # keys
    openrouter_key = os.getenv("OPENROUTER_API_KEY")
    groq_key = os.getenv("GROQ_API_KEY")

    # keep memory to avoid repetition
    used_titles = set()
    used_bullets = set()

    generated_slides = []

    slide_num = 1
    while slide_num <= slides_requested:
        batch_count = min(BATCH_SIZE, slides_requested - slide_num + 1)

        # build prompt with history
        history_block = ""
        if used_titles:
            history_block += "\nPreviously used titles:\n" + "\n".join(list(used_titles))
        if used_bullets:
            history_block += "\nPreviously used bullet ideas:\n" + "\n".join(list(used_bullets)[:200])

        prompt_text = BASE_BATCH_PROMPT.format(
            count=batch_count,
            topic=topic,
            start_index=slide_num,
            description=(description + "\n" + history_block) if history_block else description
        )

        messages = [{"role": "user", "content": prompt_text}]

        try:
            # prefer an OpenRouter free-ish model name you used earlier
            prefer_open_model = "openai/gpt-oss-120b:free"
            raw_data = call_model_with_fallback(openrouter_key, groq_key, prefer_open_model, GROQ_MODEL_ORDER, messages)
        except Exception as e:
            current_app.logger.exception("All model calls failed")
            return jsonify({"error": "Model calls failed", "details": str(e)}), 502

        # extract text (support both OpenRouter/GROQ shapes)
        try:
            text = None
            if isinstance(raw_data, dict):
                # many providers put content at choices[0].message.content
                choices = raw_data.get("choices") or raw_data.get("response") or []
                if isinstance(choices, list) and len(choices) > 0:
                    # some return choices[0]["message"]["content"], some choices[0]["text"]
                    ch = choices[0]
                    if isinstance(ch, dict):
                        msg = ch.get("message") or ch
                        text = msg.get("content") or msg.get("text") or None
                    else:
                        text = str(ch)
            if not text:
                text = str(raw_data)
        except Exception:
            text = str(raw_data)

        # parse slides
        slides_batch = parse_slides(text)
        if not slides_batch:
            current_app.logger.warning("Parser found no slides, raw model output: %s", text[:500])
            return jsonify({"error": "Model did not return slides in expected format", "raw": text}), 502

        # append parsed slides (ensuring we don't exceed requested number)
        for s in slides_batch:
            if slide_num > slides_requested:
                break
            # Simple dedupe — if title already used, skip
            title_lower = s["title"].strip().lower()
            if title_lower in used_titles:
                continue
            generated_slides.append(s)
            used_titles.add(title_lower)
            for b in s["bullets"]:
                used_bullets.add(b.strip().lower())
            slide_num += 1

        # small pause to reduce rate-hits
        time.sleep(0.2)

    # build PPT bytes
    ppt_bytes_io = build_ppt_bytes(topic, generated_slides, config={"subtitle": description})

    # Create attachment name
    safe_name = topic.strip().replace(" ", "_")[:120]
    filename = f"{safe_name}.pptx"

    # Return file as attachment
    return send_file(
        ppt_bytes_io,
        as_attachment=True,
        download_name=filename,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
